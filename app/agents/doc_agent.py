import os
import io
import base64
import json
import csv
from pathlib import Path
from typing import Dict, Any, Optional

class FileExtractor:
    """Extract text and metadata from uploaded document files."""

    @staticmethod
    def decode_file_bytes(base64_data: str) -> bytes:
        """Decode base64 file data string."""
        if "," in base64_data:
            base64_data = base64_data.split(",")[1]
        return base64.b64decode(base64_data)

    @classmethod
    def extract_text(cls, filename: str, base64_data: str) -> Dict[str, Any]:
        """Extract plain text based on file extension."""
        ext = Path(filename).suffix.lower()
        file_bytes = cls.decode_file_bytes(base64_data)
        
        extracted_text = ""
        metadata = {"filename": filename, "extension": ext, "size_bytes": len(file_bytes)}

        try:
            if ext == ".pdf":
                extracted_text = cls._extract_pdf(file_bytes)
            elif ext in [".docx", ".doc"]:
                extracted_text = cls._extract_docx(file_bytes)
            elif ext in [".pptx", ".ppt"]:
                extracted_text = cls._extract_pptx(file_bytes)
            elif ext in [".csv", ".tsv"]:
                extracted_text = cls._extract_csv(file_bytes)
            elif ext == ".json":
                extracted_text = cls._extract_json(file_bytes)
            else:
                # Default UTF-8 text decoding for txt, md, py, js, html, css, etc.
                extracted_text = file_bytes.decode("utf-8", errors="ignore")
        except Exception as e:
            extracted_text = f"[Error reading file {filename}: {str(e)}]"

        metadata["character_count"] = len(extracted_text)
        return {
            "metadata": metadata,
            "content": extracted_text.strip()
        }

    @staticmethod
    def _extract_pdf(file_bytes: bytes) -> str:
        try:
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(file_bytes))
            pages = []
            for i, page in enumerate(reader.pages):
                text = page.extract_text()
                if text:
                    pages.append(f"--- Page {i+1} ---\n{text}")
            return "\n\n".join(pages)
        except Exception:
            return file_bytes.decode("utf-8", errors="ignore")

    @staticmethod
    def _extract_docx(file_bytes: bytes) -> str:
        try:
            import docx
            doc = docx.Document(io.BytesIO(file_bytes))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n".join(paragraphs)
        except Exception:
            return file_bytes.decode("utf-8", errors="ignore")

    @staticmethod
    def _extract_pptx(file_bytes: bytes) -> str:
        try:
            import pptx
            prs = pptx.Presentation(io.BytesIO(file_bytes))
            slides_text = []
            for i, slide in enumerate(prs.slides):
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                if slide_content:
                    slides_text.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_content))
            return "\n\n".join(slides_text)
        except Exception:
            return file_bytes.decode("utf-8", errors="ignore")

    @staticmethod
    def _extract_csv(file_bytes: bytes) -> str:
        text = file_bytes.decode("utf-8", errors="ignore")
        reader = csv.reader(io.StringIO(text))
        rows = list(reader)
        if not rows:
            return text
        formatted = [f"CSV Header / Columns: {', '.join(rows[0])}"]
        for i, row in enumerate(rows[1:50]):  # Sample first 50 rows
            formatted.append(f"Row {i+1}: {', '.join(row)}")
        if len(rows) > 51:
            formatted.append(f"... total {len(rows)-1} data rows.")
        return "\n".join(formatted)

    @staticmethod
    def _extract_json(file_bytes: bytes) -> str:
        text = file_bytes.decode("utf-8", errors="ignore")
        try:
            data = json.loads(text)
            return json.dumps(data, indent=2)
        except Exception:
            return text
